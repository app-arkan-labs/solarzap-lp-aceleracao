from pptx import Presentation
import sys

def extract_text(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "--- SLIDE BREAK ---\n"
        return text
    except Exception as e:
        return str(e)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <pptx_path>")
    else:
        print(extract_text(sys.argv[1]))
